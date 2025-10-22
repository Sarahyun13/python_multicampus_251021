import os
import sys
import tempfile
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import argparse

TMDB_IMAGE_BASE = "https://image.tmdb.org/t/p/w500"

def fetch_tmdb_popular(api_key, count=8):
    items = []
    page = 1
    per_page = 20
    while len(items) < count:
        url = "https://api.themoviedb.org/3/discover/tv"
        params = {
            "api_key": api_key,
            "with_original_language": "ko",
            "sort_by": "popularity.desc",
            "page": page
        }
        r = requests.get(url, params=params, timeout=10)
        r.raise_for_status()
        data = r.json()
        for tv in data.get("results", []):
            if len(items) >= count:
                break
            items.append({
                "title": tv.get("name"),
                "overview": tv.get("overview") or "",
                "poster": (TMDB_IMAGE_BASE + tv["poster_path"]) if tv.get("poster_path") else None
            })
        if page >= data.get("total_pages", 1):
            break
        page += 1
    return items

def fetch_wikipedia_extract(title):
    S = requests.Session()
    params = {
        "action": "query",
        "prop": "extracts",
        "exintro": True,
        "explaintext": True,
        "format": "json",
        "titles": title
    }
    r = S.get("https://en.wikipedia.org/w/api.php", params=params, timeout=10)
    r.raise_for_status()
    data = r.json()
    pages = data.get("query", {}).get("pages", {})
    for p in pages.values():
        return p.get("extract", "")
    return ""

def download_image_to_file(url):
    if not url:
        return None
    r = requests.get(url, stream=True, timeout=10)
    if r.status_code != 200:
        return None
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg")
    with open(tmp.name, "wb") as f:
        for chunk in r.iter_content(1024):
            f.write(chunk)
    # ensure readable by PIL (convert if needed)
    try:
        img = Image.open(tmp.name)
        rgb = img.convert("RGB")
        rgb.save(tmp.name, format="JPEG")
    except Exception:
        pass
    return tmp.name

def add_title_slide(prs, title_text):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = title_text
    subtitle.text = "자동 생성된 K-Drama 요약 자료"
    return slide

def add_show_slide(prs, title, overview, image_path):
    blank = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank)
    left = Inches(0.5)
    top = Inches(1)
    pic_width = Inches(3.5)
    if image_path and os.path.exists(image_path):
        slide.shapes.add_picture(image_path, left, top, width=pic_width)
    # 텍스트 상자
    tx_left = left + pic_width + Inches(0.4)
    tx_top = top
    tx_width = Inches(6.0) - pic_width
    tx_height = Inches(4)
    tx = slide.shapes.add_textbox(tx_left, tx_top, tx_width, tx_height)
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = overview or "(요약 없음)"
    p.font.size = Pt(14)

def build_presentation(shows, out_path):
    prs = Presentation()
    add_title_slide(prs, "K-Drama Collection")
    temp_files = []
    try:
        for s in shows:
            img_file = None
            if s.get("poster"):
                img_file = download_image_to_file(s["poster"])
                if img_file:
                    temp_files.append(img_file)
            add_show_slide(prs, s.get("title", "제목 없음"), s.get("overview", ""), img_file)
        prs.save(out_path)
    finally:
        for f in temp_files:
            try:
                os.remove(f)
            except Exception:
                pass

def main():
    parser = argparse.ArgumentParser(description="Generate k-drama.pptx")
    parser.add_argument("--tmdb-key", help="TMDB API Key (optional)")
    parser.add_argument("--count", type=int, default=8, help="몇 편 수집할지")
    parser.add_argument("--out", default=os.path.join(os.getcwd(), "k-drama.pptx"), help="출력 파일 경로")
    args = parser.parse_args()

    shows = []
    if args.tmdb_key:
        try:
            shows = fetch_tmdb_popular(args.tmdb_key, count=args.count)
        except Exception as e:
            print("TMDB API 오류:", e, file=sys.stderr)

    if not shows:
        # TMDB 키 없거나 실패 시 기본 목록 + 위키백과 요약
        default_titles = [
            "Squid Game", "Crash Landing on You", "Goblin (TV series)",
            "Descendants of the Sun", "Itaewon Class", "Reply 1988",
            "My Mister (TV series)", "Mr. Sunshine (TV series)"
        ][:args.count]
        for t in default_titles:
            overview = fetch_wikipedia_extract(t)
            shows.append({"title": t.replace(" (TV series)", ""), "overview": overview, "poster": None})

    out_path = os.path.abspath(args.out)
    build_presentation(shows, out_path)
    print("생성 완료:", out_path)

if __name__ == "__main__":
    main()